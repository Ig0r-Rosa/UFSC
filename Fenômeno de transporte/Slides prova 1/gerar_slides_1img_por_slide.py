from __future__ import annotations

import os
from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt


@dataclass(frozen=True)
class Layout:
    slide_w_in: float = 13.333  # 16:9
    slide_h_in: float = 7.5
    margin_in: float = 0.35
    title_h_in: float = 0.55


def list_images(folder: Path) -> list[Path]:
    exts = {".png", ".jpg", ".jpeg", ".webp"}
    imgs = [p for p in folder.iterdir() if p.is_file() and p.suffix.lower() in exts]
    return sorted(imgs, key=lambda p: p.name)


def fit_aspect(iw: int, ih: int, bw: float, bh: float) -> tuple[float, float]:
    if iw <= 0 or ih <= 0 or bw <= 0 or bh <= 0:
        return bw, bh
    ir = iw / ih
    br = bw / bh
    if ir >= br:
        w = bw
        h = bw / ir
    else:
        h = bh
        w = bh * ir
    return w, h


def add_title(slide, text: str, layout: Layout) -> None:
    box = slide.shapes.add_textbox(
        Inches(layout.margin_in),
        Inches(layout.margin_in * 0.6),
        Inches(layout.slide_w_in - 2 * layout.margin_in),
        Inches(layout.title_h_in),
    )
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(24)


def add_image_slide(prs: Presentation, img_path: Path, idx: int, total: int, folder_name: str, layout: Layout) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_title(slide, f"{folder_name} — {idx}/{total}", layout)

    avail_w = layout.slide_w_in - 2 * layout.margin_in
    top = layout.margin_in + layout.title_h_in
    avail_h = layout.slide_h_in - top - layout.margin_in

    try:
        with Image.open(img_path) as im:
            iw, ih = im.size
    except Exception:
        iw, ih = (1, 1)

    w, h = fit_aspect(iw, ih, avail_w, avail_h)
    left = layout.margin_in + (avail_w - w) / 2
    upper = top + (avail_h - h) / 2

    slide.shapes.add_picture(
        str(img_path),
        Inches(left),
        Inches(upper),
        width=Inches(w),
        height=Inches(h),
    )


def build_deck_for_folder(folder: Path, out_path: Path, layout: Layout) -> None:
    images = list_images(folder)
    prs = Presentation()
    prs.slide_width = Inches(layout.slide_w_in)
    prs.slide_height = Inches(layout.slide_h_in)

    total = len(images)
    for i, img in enumerate(images, start=1):
        add_image_slide(prs, img, i, total, folder.name, layout)

    prs.save(out_path)


def main() -> int:
    base = Path(__file__).resolve().parent
    layout = Layout()

    for i in range(1, 5):
        folder = base / f"Slide {i}"
        if not folder.exists():
            raise FileNotFoundError(f"Pasta não encontrada: {folder}")
        out = base / f"Slides-Prova1-{folder.name.replace(' ', '')}-1imgPorSlide.pptx"
        build_deck_for_folder(folder, out, layout)
        print(f"Gerado: {out}")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())

