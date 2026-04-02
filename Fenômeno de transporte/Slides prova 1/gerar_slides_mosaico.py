from __future__ import annotations

import math
import os
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


@dataclass(frozen=True)
class Layout:
    margin_in: float = 0.4
    title_h_in: float = 0.55
    gap_in: float = 0.06


def list_images(folder: Path) -> list[Path]:
    exts = {".png", ".jpg", ".jpeg", ".webp"}
    imgs = [p for p in folder.iterdir() if p.is_file() and p.suffix.lower() in exts]
    return sorted(imgs, key=lambda p: p.name)


def best_grid(n: int) -> tuple[int, int]:
    if n <= 0:
        return (1, 1)
    cols = max(1, math.ceil(math.sqrt(n)))
    rows = math.ceil(n / cols)
    return cols, rows


def fit_aspect(
    img_w_px: int,
    img_h_px: int,
    box_w_in: float,
    box_h_in: float,
) -> tuple[float, float]:
    if img_w_px <= 0 or img_h_px <= 0:
        return box_w_in, box_h_in
    img_ratio = img_w_px / img_h_px
    box_ratio = box_w_in / box_h_in if box_h_in else img_ratio

    if img_ratio >= box_ratio:
        w = box_w_in
        h = box_w_in / img_ratio
    else:
        h = box_h_in
        w = box_h_in * img_ratio
    return w, h


def chunked(it: Iterable[Path], size: int) -> list[list[Path]]:
    out: list[list[Path]] = []
    buf: list[Path] = []
    for item in it:
        buf.append(item)
        if len(buf) >= size:
            out.append(buf)
            buf = []
    if buf:
        out.append(buf)
    return out


def add_mosaic_slide(
    prs: Presentation,
    title: str,
    images: list[Path],
    layout: Layout,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title-only
    slide.shapes.title.text = f"{title} ({len(images)} imagens)"

    sw = prs.slide_width
    sh = prs.slide_height
    sw_in = sw / 914400  # EMU per inch
    sh_in = sh / 914400

    left_in = layout.margin_in
    top_in = layout.margin_in + layout.title_h_in
    avail_w_in = sw_in - 2 * layout.margin_in
    avail_h_in = sh_in - top_in - layout.margin_in

    cols, rows = best_grid(len(images))
    cell_w_in = (avail_w_in - (cols - 1) * layout.gap_in) / cols
    cell_h_in = (avail_h_in - (rows - 1) * layout.gap_in) / rows

    for idx, img_path in enumerate(images):
        r = idx // cols
        c = idx % cols
        x_in = left_in + c * (cell_w_in + layout.gap_in)
        y_in = top_in + r * (cell_h_in + layout.gap_in)

        try:
            with Image.open(img_path) as im:
                iw, ih = im.size
        except Exception:
            iw, ih = (cell_w_in, cell_h_in)  # fallback

        w_in, h_in = fit_aspect(iw, ih, cell_w_in, cell_h_in)
        dx_in = x_in + (cell_w_in - w_in) / 2
        dy_in = y_in + (cell_h_in - h_in) / 2

        slide.shapes.add_picture(
            str(img_path),
            Inches(dx_in),
            Inches(dy_in),
            width=Inches(w_in),
            height=Inches(h_in),
        )


def main() -> int:
    base = Path(__file__).resolve().parent
    folders = [base / f"Slide {i}" for i in range(1, 5)]
    missing = [str(p) for p in folders if not p.exists()]
    if missing:
        raise FileNotFoundError(f"Pastas não encontradas: {', '.join(missing)}")

    prs = Presentation()
    # widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    layout = Layout()

    # Limite prático: muita imagem em 1 slide pode ficar minúscula.
    # Mantemos 1 slide por pasta, como pedido.
    for folder in folders:
        imgs = list_images(folder)
        add_mosaic_slide(prs, folder.name, imgs, layout)

    out_path = base / "Slides-Prova1-mosaico.pptx"
    prs.save(out_path)
    print(f"Gerado: {out_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

